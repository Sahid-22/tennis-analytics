import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Create presentation
prs = Presentation()

# Set 16:9 aspect ratio
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BG_COLOR = RGBColor(10, 10, 26)       # #0A0A1A
CARD_COLOR = RGBColor(17, 17, 40)     # #111128 (approx)
TEXT_COLOR = RGBColor(226, 232, 240)  # #E2E8F0
CYAN = RGBColor(0, 212, 255)          # #00D4FF
PURPLE = RGBColor(124, 58, 237)       # #7C3AED
GREEN = RGBColor(16, 185, 129)        # #10B981
AMBER = RGBColor(245, 158, 11)       # #F59E0B
ROSE = RGBColor(244, 63, 94)          # #F43F5E

# Helper functions
def set_slide_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_accent_line(slide, x, y, width, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(0.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_slide_number(slide, number):
    txBox = slide.shapes.add_textbox(Inches(12.5), Inches(6.8), Inches(0.5), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.name = 'Calibri'
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_COLOR
    p.alignment = PP_ALIGN.RIGHT

def set_text_format(run, size=18, bold=False, color=TEXT_COLOR):
    run.font.name = 'Calibri'
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color

def add_title(slide, text):
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    set_text_format(p.runs[0], size=36, bold=True)
    add_accent_line(slide, 0.8, 1.4, 2, CYAN)

def add_metric_card(slide, x, y, width, height, value, label, color=CYAN):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_COLOR
    shape.line.color.rgb = color
    shape.line.width = Pt(1.5)
    
    # Value
    txBox = slide.shapes.add_textbox(Inches(x+0.1), Inches(y+height/2-0.5), Inches(width-0.2), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = value
    p.alignment = PP_ALIGN.CENTER
    set_text_format(p.runs[0], size=32, bold=True, color=color)
    
    # Label
    txBox2 = slide.shapes.add_textbox(Inches(x+0.1), Inches(y+height/2+0.1), Inches(width-0.2), Inches(0.4))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = label
    p2.alignment = PP_ALIGN.CENTER
    set_text_format(p2.runs[0], size=14)

def add_card(slide, x, y, width, height, title, items, color=CYAN):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_COLOR
    shape.line.color.rgb = color
    shape.line.width = Pt(1.5)
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(x+0.2), Inches(y+0.2), Inches(width-0.4), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    set_text_format(p.runs[0], size=20, bold=True, color=color)
    
    # Items
    txBox2 = slide.shapes.add_textbox(Inches(x+0.2), Inches(y+0.7), Inches(width-0.4), Inches(height-0.9))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for item in items:
        p2 = tf2.add_paragraph()
        p2.text = "• " + item
        set_text_format(p2.runs[0], size=16)

def add_basic_list(slide, x, y, items):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(11.7), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for item in items:
        p = tf.add_paragraph()
        p.text = "• " + item
        set_text_format(p.runs[0], size=20)
        p.space_after = Pt(10)

blank_slide_layout = prs.slide_layouts[6]

# SLIDE 1: Title Slide
s1 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(s1)
tx1 = s1.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(8), Inches(1.5))
p1 = tx1.text_frame.paragraphs[0]
p1.text = "Tennis Analytics"
set_text_format(p1.runs[0], size=54, bold=True, color=CYAN)
p2 = tx1.text_frame.add_paragraph()
p2.text = "A Professional Sports Data Analytics Platform"
set_text_format(p2.runs[0], size=24, color=TEXT_COLOR)

add_metric_card(s1, 9, 2, 3.5, 1.2, "6,619", "Competitions", CYAN)
add_metric_card(s1, 9, 3.4, 3.5, 1.2, "1,000", "Competitors", PURPLE)
add_metric_card(s1, 9, 4.8, 3.5, 1.2, "4,021", "Venues", GREEN)

tx2 = s1.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.5))
p3 = tx2.text_frame.paragraphs[0]
p3.text = "Powered by SportRadar API • Python • Streamlit • Plotly | Version 2.0 | July 2026"
set_text_format(p3.runs[0], size=14, color=TEXT_COLOR)
add_slide_number(s1, 1)

# SLIDE 2: Table of Contents
s2 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(s2)
add_title(s2, "Table of Contents")
add_card(s2, 0.8, 2, 5.5, 4.5, "Part I: Architecture & Data", ["Project Overview", "Problem & Objectives", "Tech Stack & Architecture", "Data Pipeline (ETL)", "API & Database Design", "Data Quality Framework"], CYAN)
add_card(s2, 6.8, 2, 5.5, 4.5, "Part II: Dashboard & DevOps", ["Dashboard Overview & Pages", "Advanced Analytics", "Testing Strategy", "CI/CD & Deployment", "Security & Metrics", "Challenges & Future Outlook"], PURPLE)
add_slide_number(s2, 2)

# SLIDE 3: Project Overview
s3 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(s3)
add_title(s3, "Project Overview")
add_basic_list(s3, 0.8, 1.8, ["A comprehensive platform for fetching, processing, and visualizing professional tennis data."])
add_metric_card(s3, 0.8, 3, 2.7, 2, "3", "API Endpoints", CYAN)
add_metric_card(s3, 3.8, 3, 2.7, 2, "7", "SQL Tables", PURPLE)
add_metric_card(s3, 6.8, 3, 2.7, 2, "7", "Dashboard Pages", GREEN)
add_metric_card(s3, 9.8, 3, 2.7, 2, "38", "Unit Tests", AMBER)
add_slide_number(s3, 3)

# SLIDE 4: Problem Statement & Objectives
s4 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(s4)
add_title(s4, "Problem Statement & Objectives")
add_card(s4, 0.8, 2, 5.5, 3.5, "Problem Statement", ["Scattered, inaccessible sports data", "Complex nested JSON structures", "No data quality checks in raw data", "Manual, error-prone processing"], ROSE)
add_card(s4, 6.8, 2, 5.5, 3.5, "Objectives", ["Automated robust ETL pipeline", "Relational database schema", "Interactive analytical dashboard", "Automated quality validation"], GREEN)
add_slide_number(s4, 4)

# SLIDE 5: Technology Stack
s5 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(s5)
add_title(s5, "Technology Stack")
add_card(s5, 0.8, 2, 2.7, 4.5, "Backend", ["Python 3.11+", "SQLAlchemy 2.0", "Requests", "Pandas"], CYAN)
add_card(s5, 3.8, 2, 2.7, 4.5, "Frontend", ["Streamlit", "Plotly", "Vanilla CSS", "Google Fonts"], PURPLE)
add_card(s5, 6.8, 2, 2.7, 4.5, "DevOps", ["pytest", "GitHub Actions", "Docker", "Ruff (Linter)"], GREEN)
add_card(s5, 9.8, 2, 2.7, 4.5, "Data", ["SQLite / PostgreSQL", "JSON / CSV", "SHA-256 Hashing", "API v3"], AMBER)
add_slide_number(s5, 5)

# SLIDE 6: System Architecture
s6 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(s6)
add_title(s6, "System Architecture")
add_basic_list(s6, 0.8, 2, ["SportRadar API \u2192 api_client.py \u2192 transformers.py \u2192 database.py \u2192 app.py", "Component support modules: snapshots, quality, logging"])
add_card(s6, 0.8, 3.5, 11.7, 3, "Workflow Components", ["Data extraction with rate limit handling", "Data normalization and flattening", "Relational loading and indexing", "Interactive UI generation"], CYAN)
add_slide_number(s6, 6)

# SLIDE 7: Data Pipeline (ETL)
s7 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(s7)
add_title(s7, "Data Pipeline (ETL)")
add_card(s7, 0.8, 2, 2.7, 4.5, "Extract", ["api_client.py", "Endpoint mapping", "Error handling", "Payload retrieval"], CYAN)
add_card(s7, 3.8, 2, 2.7, 4.5, "Transform", ["transformers.py", "JSON flattening", "Type casting", "Null handling"], PURPLE)
add_card(s7, 6.8, 2, 2.7, 4.5, "Load", ["database.py", "Schema enforcing", "Upsert operations", "Foreign Keys"], GREEN)
add_card(s7, 9.8, 2, 2.7, 4.5, "Validate", ["quality.py", "Data freshness", "Completeness", "Integrity checks"], AMBER)
add_slide_number(s7, 7)

# SLIDE 8: API Integration
s8 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(s8)
add_title(s8, "API Integration")
add_card(s8, 0.8, 2, 11.7, 1.2, "Competitions Endpoint", ["URL: /tournaments/en/competitions.json", "6,619 Records fetched"], CYAN)
add_card(s8, 0.8, 3.5, 11.7, 1.2, "Complexes Endpoint", ["URL: /tournaments/en/complexes.json", "767 Records fetched"], PURPLE)
add_card(s8, 0.8, 5, 11.7, 1.2, "Doubles Rankings", ["URL: /rankings/en/competitors/doubles.json", "1,000 Records fetched"], GREEN)
add_slide_number(s8, 8)

# SLIDE 9: Database Design
s9 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(s9)
add_title(s9, "Database Design")
add_card(s9, 0.8, 2, 11.7, 4.5, "Relational Schema (7 Tables, ~13,428 Rows)", ["categories (18 rows)", "competitions (6,619 rows)", "complexes (767 rows)", "venues (4,021 rows)", "competitors (1,000 rows)", "competitor_rankings (1,000 rows)", "api_sync_log (3 rows)"], CYAN)
add_slide_number(s9, 9)

# SLIDE 10: Data Transformation
s10 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(s10)
add_title(s10, "Data Transformation")
add_card(s10, 0.8, 2, 5.5, 4.5, "Raw API JSON", ["Nested dictionaries", "Array of objects", "Mixed data types", "Inconsistent nulls"], ROSE)
add_card(s10, 6.8, 2, 5.5, 4.5, "Transformation Steps", ["Flatten nested structures", "Standardize date formats", "Type enforce integers/floats", "Handle missing optional keys", "Normalize foreign keys", "Extract sub-lists to tables", "Clean text fields", "Generate synthetic keys if needed"], GREEN)
add_slide_number(s10, 10)

# SLIDE 11: Data Quality Framework
s11 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(s11)
add_title(s11, "Data Quality Framework")
add_card(s11, 0.8, 2, 5.5, 2, "Completeness", ["Checks for non-null required fields", "5 checks implemented"], CYAN)
add_card(s11, 6.8, 2, 5.5, 2, "Integrity", ["Validates foreign key relationships", "4 checks implemented"], PURPLE)
add_card(s11, 0.8, 4.3, 5.5, 2, "Validity", ["Checks data ranges and formats", "3 checks implemented"], GREEN)
add_card(s11, 6.8, 4.3, 5.5, 2, "Freshness", ["Verifies recent API sync dates", "3 checks implemented"], AMBER)
add_slide_number(s11, 11)

# SLIDE 12: Dashboard Design Overview
s12 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(s12)
add_title(s12, "Dashboard Design Overview")
add_card(s12, 0.8, 2, 2.7, 2, "Dark Theme", ["Glassmorphism", "Premium Look"], CYAN)
add_card(s12, 3.8, 2, 2.7, 2, "Interactive", ["Plotly Charts", "Hover Effects"], PURPLE)
add_card(s12, 6.8, 2, 2.7, 2, "Animations", ["Micro-interactions", "CSS Transitions"], GREEN)
add_card(s12, 9.8, 2, 2.7, 2, "Export", ["CSV Download", "JSON Export"], AMBER)
add_card(s12, 0.8, 4.3, 11.7, 2.2, "7 Pages", ["Overview, Competitions, Competitors, Venues, SQL Analysis, Advanced Analytics, Data Quality"], CYAN)
add_slide_number(s12, 12)

# SLIDE 13: Dashboard - Overview Page
s13 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(s13)
add_title(s13, "Dashboard — Overview Page")
add_basic_list(s13, 0.8, 2, ["High-level system metrics", "Animated gradient metric cards", "Quick navigation links", "Recent sync status", "Data volume sparklines", "System health indicator", "API quota usage"])
add_slide_number(s13, 13)

# SLIDE 14: Dashboard - Competitions Page
s14 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(s14)
add_title(s14, "Dashboard — Competitions Page")
add_basic_list(s14, 0.8, 2, ["Grid view of 6,619 competitions", "Filter by category/gender", "Search by tournament name", "Detailed competition modal", "Timeline of events", "Category distribution pie chart", "Export functionality"])
add_slide_number(s14, 14)

# SLIDE 15: Dashboard - Competitors Page
s15 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(s15)
add_title(s15, "Dashboard — Competitors Page")
add_basic_list(s15, 0.8, 2, ["List of 1,000 top ranked competitors", "Doubles team analysis", "Points progression charts", "Country distribution map", "Age/Experience demographics", "Head-to-head comparison tool", "Performance metrics"])
add_slide_number(s15, 15)

# SLIDE 16: Dashboard - Venues & SQL Pages
s16 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(s16)
add_title(s16, "Dashboard — Venues & SQL Pages")
add_basic_list(s16, 0.8, 2, ["Venues: Geographic mapping of 4,021 venues", "Venues: Capacity analytics", "Venues: Surface type breakdown", "Venues: Complex association", "SQL: Custom interactive query editor", "SQL: Schema reference browser", "SQL: Query result export"])
add_slide_number(s16, 16)

# SLIDE 17: Dashboard - Advanced Analytics
s17 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(s17)
add_title(s17, "Dashboard — Advanced Analytics")
add_basic_list(s17, 0.8, 2, ["Complex interactive visualizations", "Multi-variable bubble charts", "Geographic choropleth maps", "Hierarchical sunburst charts", "Correlation heatmaps", "Trend forecasting models", "Dynamic cross-filtering"])
add_slide_number(s17, 17)

# SLIDE 18: Dashboard - Data Quality Page
s18 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(s18)
add_title(s18, "Dashboard — Data Quality Page")
add_basic_list(s18, 0.8, 2, ["Composite quality score (0-100)", "Dimension radar chart", "Failed check detailed logs", "Historical quality trends", "Automated anomaly detection", "Missing value heatmap", "Schema validation results"])
add_slide_number(s18, 18)

# SLIDE 19: Key Features Summary
s19 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(s19)
add_title(s19, "Key Features Summary")
add_card(s19, 0.8, 2, 5.5, 2, "ETL Pipeline & Quality", ["Robust data extraction and validation"], CYAN)
add_card(s19, 6.8, 2, 5.5, 2, "Dark Theme & UI", ["Premium glassmorphism aesthetics"], PURPLE)
add_card(s19, 0.8, 4.3, 5.5, 2, "Interactive Charts", ["Plotly-driven dynamic visuals"], GREEN)
add_card(s19, 6.8, 4.3, 5.5, 2, "DevOps Ready", ["Dockerized with CI/CD GitHub Actions"], AMBER)
add_slide_number(s19, 19)

# SLIDE 20: Testing Strategy
s20 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(s20)
add_title(s20, "Testing Strategy")
add_metric_card(s20, 0.8, 2, 4, 2, "38", "Tests Passing in 1.77s", GREEN)
add_card(s20, 5.2, 2, 7.3, 4.5, "Test Modules (8 Total)", ["test_api_client.py - API mocking", "test_transformers.py - Logic verify", "test_database.py - Schema checks", "test_quality.py - Validation rules", "test_app.py - UI components", "test_integration.py - E2E flow", "conftest.py - Fixtures", "Focus on 90%+ code coverage"], CYAN)
add_slide_number(s20, 20)

# SLIDE 21: CI/CD Pipeline
s21 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(s21)
add_title(s21, "CI/CD Pipeline")
add_basic_list(s21, 0.8, 2, ["Triggered on Push / Pull Request to main", "Stages: Checkout \u2192 Setup Python \u2192 Install Deps", "Stage: Lint (Ruff)", "Stage: Typecheck (Mypy)", "Stage: Test (pytest)", "Pre-commit hooks for local validation"])
add_slide_number(s21, 21)

# SLIDE 22: Deployment (Docker)
s22 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(s22)
add_title(s22, "Deployment (Docker)")
add_card(s22, 0.8, 2, 5.5, 2, "Multi-Stage Build", ["Optimized image size"], CYAN)
add_card(s22, 6.8, 2, 5.5, 2, "Non-Root User", ["Enhanced container security"], PURPLE)
add_card(s22, 0.8, 4.3, 5.5, 2, "Health Check", ["Automatic container monitoring"], GREEN)
add_card(s22, 6.8, 4.3, 5.5, 2, "Layer Caching", ["Faster build times"], AMBER)
add_slide_number(s22, 22)

# SLIDE 23: Security Best Practices
s23 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(s23)
add_title(s23, "Security Best Practices")
add_card(s23, 0.8, 2, 11.7, 1.2, "API Security", ["Key loaded from .env", "No hardcoded credentials", "Rate limiting handled gracefully"], CYAN)
add_card(s23, 0.8, 3.5, 11.7, 1.2, "Database Security", ["SQL Injection prevention via SQLAlchemy", "Sanitized inputs", "Secure file permissions"], PURPLE)
add_card(s23, 0.8, 5, 11.7, 1.2, "Application Security", ["Dependency scanning", "Docker non-root execution", "Security headers in UI"], GREEN)
add_slide_number(s23, 23)

# SLIDE 24: Results & Metrics
s24 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(s24)
add_title(s24, "Results & Metrics")
add_metric_card(s24, 0.8, 2, 3.5, 1.2, "13,428", "Records", CYAN)
add_metric_card(s24, 4.9, 2, 3.5, 1.2, "7", "Tables", PURPLE)
add_metric_card(s24, 9.0, 2, 3.5, 1.2, "15+", "Quality Checks", GREEN)
add_metric_card(s24, 0.8, 3.5, 3.5, 1.2, "38", "Tests", AMBER)
add_metric_card(s24, 4.9, 3.5, 3.5, 1.2, "7", "Pages", ROSE)
add_metric_card(s24, 9.0, 3.5, 3.5, 1.2, "35+", "Source Files", CYAN)
add_card(s24, 0.8, 5.2, 11.7, 1.5, "Before vs After", ["Manual scripts -> Automated ETL", "No UI -> Interactive 7-page Dashboard", "Fragile parsing -> Validated schema constraints"], GREEN)
add_slide_number(s24, 24)

# SLIDE 25: Challenges & Solutions
s25 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(s25)
add_title(s25, "Challenges & Solutions")
add_basic_list(s25, 0.8, 2, [
    "Challenge: Nested JSON -> Solution: Custom transformers.py mapping",
    "Challenge: API Rate Limits -> Solution: Backoff & retry mechanism",
    "Challenge: Multi-DB -> Solution: SQLAlchemy ORM abstraction",
    "Challenge: Dashboard Lag -> Solution: Streamlit caching (@st.cache_data)",
    "Challenge: Test Coverage -> Solution: Pytest fixtures & mock responses"
])
add_slide_number(s25, 25)

# SLIDE 26: Future Enhancements
s26 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(s26)
add_title(s26, "Future Enhancements")
add_card(s26, 0.8, 2, 3.7, 4.5, "Phase 1: Near-Term", ["Live match data", "Player headshots", "More API endpoints", "Performance tuning"], CYAN)
add_card(s26, 4.8, 2, 3.7, 4.5, "Phase 2: Mid-Term", ["Predictive ML models", "User authentication", "Custom alert emails", "PostgreSQL migration"], PURPLE)
add_card(s26, 8.8, 2, 3.7, 4.5, "Phase 3: Long-Term", ["Mobile application", "Real-time WebSockets", "Public API Gateway", "Multi-sport expansion"], GREEN)
add_slide_number(s26, 26)

# SLIDE 27: Conclusion
s27 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(s27)
add_title(s27, "Conclusion")
add_card(s27, 0.8, 2, 11.7, 4.5, "Project Impact", [
    "Successfully built a robust end-to-end data platform",
    "Mastered API integration and complex data transformation",
    "Implemented professional-grade database schemas",
    "Created an engaging, premium user interface",
    "Applied strict software engineering practices (Testing, CI/CD, Docker)",
    "Ready for scalable production deployment"
], CYAN)
add_slide_number(s27, 27)

# SLIDE 28: References
s28 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(s28)
add_title(s28, "References")
add_basic_list(s28, 0.8, 2, [
    "SportRadar Tennis API Docs: developer.sportradar.com",
    "Python 3.11 Documentation: docs.python.org/3",
    "Streamlit API Reference: docs.streamlit.io",
    "SQLAlchemy 2.0 Docs: docs.sqlalchemy.org",
    "Plotly Graphing Library: plotly.com/python",
    "Pytest Documentation: docs.pytest.org",
    "Docker Multi-stage Builds: docs.docker.com",
    "GitHub Actions: docs.github.com/en/actions",
    "Pandas Documentation: pandas.pydata.org",
    "python-pptx: python-pptx.readthedocs.io"
])
add_slide_number(s28, 28)

# SLIDE 29: Thank You
s29 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(s29)
tx3 = s29.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.5))
p4 = tx3.text_frame.paragraphs[0]
p4.text = "Thank You"
p4.alignment = PP_ALIGN.CENTER
set_text_format(p4.runs[0], size=64, bold=True, color=CYAN)

tx4 = s29.shapes.add_textbox(Inches(0.8), Inches(4), Inches(11.7), Inches(1))
p5 = tx4.text_frame.paragraphs[0]
p5.text = "Tennis Analytics Platform • Python • Streamlit • Plotly"
p5.alignment = PP_ALIGN.CENTER
set_text_format(p5.runs[0], size=24, color=TEXT_COLOR)

tx5 = s29.shapes.add_textbox(Inches(0.8), Inches(5), Inches(11.7), Inches(1))
p6 = tx5.text_frame.paragraphs[0]
p6.text = "Questions & Feedback Welcome"
p6.alignment = PP_ALIGN.CENTER
set_text_format(p6.runs[0], size=20, color=PURPLE)
add_slide_number(s29, 29)

# Save presentation
out_path = os.path.join(r"d:\Labmentix Project\SportRadar_Tennis_Analytics", "Tennis_Analytics_Presentation.pptx")
prs.save(out_path)
print(f"Presentation generated successfully at: {out_path}")
print(f"Total slides: {len(prs.slides)}")
